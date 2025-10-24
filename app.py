from flask import Flask, request, jsonify, render_template, send_file, url_for, g
from werkzeug.utils import secure_filename
import os
import fitz
from docx import Document
import logging
import re
import json
import hashlib 
import time 
from typing import Dict, List, Tuple, Union
from collections import Counter
import requests
from pptx import Presentation

# NLP Libraries
import spacy
from sklearn.feature_extraction.text import TfidfVectorizer
from nltk.tokenize import sent_tokenize
import nltk

# Vector DB & Embeddings (Chroma Cloud)
import chromadb
from chromadb.api import ClientAPI
from chromadb.api.models.Collection import Collection
from chromadb import Settings
from dotenv import load_dotenv

# Rate Limiting
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address

# Load environment variables from .env file
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Download required NLTK data
try:
    nltk.download('punkt', quiet=True)
    nltk.download('stopwords', quiet=True)
    logger.info("NLTK 'punkt' and 'stopwords' downloaded.")
except Exception as e:
    logger.warning(f"Initial NLTK download failed. Sentence tokenization may rely on fallback. Error: {e}")
    pass

UPLOAD_FOLDER = "uploads"
ALLOWED_EXTENSIONS = {'pdf', 'txt', 'docx', 'ppt', 'pptx'} 
MAX_CONTENT_LIMIT = 100 * 1024 * 1024 
MAX_FILE_SIZE = 100 * 1024 * 1024  # 10MB (This is the logic limit)

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = MAX_CONTENT_LIMIT
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# --- Configuration for Security and Caching ---
app.config['SECRET_KEY'] = os.environ.get('FLASK_SECRET_KEY', 'default-dev-secret-key')
API_KEY = os.environ.get('MW_API_KEY', '1234354') 
FILE_CACHE: Dict[str, Dict] = {} 
# --- END Configuration ---

# --- Rate Limiting Initialization ---
limiter = Limiter(
    get_remote_address,
    app=app,
    default_limits=["200 per day", "50 per hour"],
    storage_uri="memory://" 
)
# --- END Rate Limiting ---

# --- NEW: JSON Handler for Rate Limits ---
@app.errorhandler(429)
def ratelimit_handler(e):
    return jsonify({"error": "Rate limit exceeded. Too many requests."}), 429
# --- END NEW ---

# Load spaCy model (lightweight)
try:
    nlp = spacy.load("en_core_web_sm")
    logger.info("spaCy model loaded")
except:
    logger.warning("spaCy model not found. Run: python -m spacy download en_core_web_sm")
    nlp = None


# ----------------------------------------------------------------------
# ChromaDB Cloud Context Management Functions (RAG Backend)
# ----------------------------------------------------------------------
def get_chroma_client() -> Union[ClientAPI, None]:
    """Initializes and returns the Chroma Cloud client, storing it in Flask's 'g'."""
    if 'chroma_client' not in g:
        try:
            g.chroma_client = chromadb.CloudClient(
                api_key=os.getenv("CHROMA_API_KEY"),
                tenant=os.getenv("CHROMA_TENANT"),
                database=os.getenv("CHROMA_DATABASE"),
                settings=Settings(allow_reset=True)
            )
            logger.info("Chroma DB Cloud Client initialized.")
        except Exception as e:
            logger.error(f"Failed to initialize Chroma Cloud Client. Check .env: {e}")
            g.chroma_client = None
    return g.chroma_client
    
def get_chroma_collection() -> Union[Collection, None]:
    """Returns the Chroma collection object, ensuring the client is initialized."""
    chroma_client = get_chroma_client()
    if not chroma_client:
        return None
        
    if 'chroma_collection' not in g:
        try:
            g.chroma_collection = chroma_client.get_or_create_collection(name="mw_tech_docs_collection")
            logger.info(f"Chroma DB Collection 'mw_tech_docs_collection' accessed. Count: {g.chroma_collection.count()}")
        except Exception as e:
            logger.error(f"Failed to access Chroma Collection: {e}")
            g.chroma_collection = None

    return g.chroma_collection

@app.teardown_appcontext
def close_chroma_client(exception=None):
    """Closes the Chroma client connection on application context teardown."""
    g.pop('chroma_client', None) 


# -------------------------
# Utility Functions
# -------------------------
def allowed_file(filename: str) -> bool:
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in {'pdf', 'txt', 'docx', 'ppt', 'pptx'}

def calculate_file_hash(file_path: str) -> str:
    """Calculates the SHA256 hash of a file's content."""
    hash_sha256 = hashlib.sha256()
    with open(file_path, 'rb') as f:
        for byte_block in iter(lambda: f.read(4096), b""):
            hash_sha256.update(byte_block)
    return hash_sha256.hexdigest()

def extract_text(file_path: str, file_type: str) -> str:
    """Extract text from PDF, DOCX, TXT, or PPTX."""
    try:
        text = ""
        if file_type == 'pdf':
            doc = fitz.open(file_path)
            for page in doc: text += page.get_text()
            doc.close()
        elif file_type == 'docx':
            doc = Document(file_path)
            text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        elif file_type == 'txt':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f: text = f.read()
        elif file_type in ['ppt', 'pptx']:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame: text += shape.text_frame.text + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting text: {e}")
        raise

def safe_sent_tokenize(text: str) -> List[str]:
    """Robust sentence tokenizer with NLTK fallback."""
    if not text: return []
    try: return sent_tokenize(text)
    except Exception as e:
        sentences = re.split(r'(?<=[.!?])\s+', text.strip())
        return [s.strip() for s in sentences if s.strip()]

def summarize_with_nlp(text: str, num_sentences: int = 5) -> str:
    """Extract key sentences using TF-IDF scoring."""
    if not text or len(text.split()) < 20: return text
    try:
        sentences = safe_sent_tokenize(text) 
        if len(sentences) <= num_sentences: return text
        vectorizer = TfidfVectorizer(max_features=100, stop_words='english')
        try: sentence_vectors = vectorizer.fit_transform(sentences)
        except: return " ".join(sentences[:num_sentences])
        sentence_scores = sentence_vectors.sum(axis=1).A1
        top_indices = sentence_scores.argsort()[-num_sentences:][::-1]
        top_indices_sorted = sorted(top_indices)
        summary_sentences = [sentences[i] for i in top_indices_sorted]
        return " ".join(summary_sentences)
    except Exception as e:
        logger.error(f"Summarization error: {e}")
        sentences = safe_sent_tokenize(text)
        return " ".join(sentences[:num_sentences])

def extract_endpoint_description(text: str, position: int, context_size: int = 200) -> str:
    start = max(0, position - context_size); end = min(len(text), position + context_size); context = text[start:end]
    sentences = safe_sent_tokenize(context) 
    if sentences: return sentences[len(sentences)//2][:100].strip()
    return "No description available"

def extract_apis_with_nlp(text: str) -> List[Dict]:
    """
    Extract API endpoints using regex patterns.
    FIXED: Regex tolerates markdown/table formatting.
    """
    apis = []
    patterns = [
        r'(GET|POST|PUT|DELETE|PATCH)\s*[^/\w]*\s*(/[\w/\-{}:]+)', # Critical fix for tables
        r'https?://[^\s]+/api/[^\s]+',
        r'(/api/v?\d*/[\w/\-]+)',
    ]
    methods = ['GET', 'POST', 'PUT', 'DELETE', 'PATCH', 'OPTIONS', 'HEAD']; found_endpoints = set()
    for pattern in patterns:
        matches = re.finditer(pattern, text, re.IGNORECASE | re.MULTILINE)
        for match in matches:
            if len(match.groups()) >= 2: method = match.group(1).upper(); endpoint = match.group(2)
            else: endpoint = match.group(0); context_start = max(0, match.start() - 50); context = text[context_start:match.start()]; method = 'GET' 
            
            endpoint = endpoint.split()[0] if ' ' in endpoint else endpoint
            endpoint = endpoint.rstrip('.,;:`\'"|') # Cleaning logic
            
            description = extract_endpoint_description(text, match.start())
            key = f"{method}:{endpoint}"
            if key not in found_endpoints and endpoint.startswith('/'): 
                found_endpoints.add(key)
                apis.append({"endpoint": endpoint, "method": method, "description": description})
    
    doc_pattern = r'(?:Endpoint|API|Route):\s*([A-Z]+)?\s*(/[\w/\-{}:]+)'; doc_matches = re.finditer(doc_pattern, text, re.IGNORECASE)
    for match in doc_matches:
        method = match.group(1).upper() if match.group(1) else 'GET'; endpoint = match.group(2); description = extract_endpoint_description(text, match.start())
        key = f"{method}:{endpoint}";
        if key not in found_endpoints: found_endpoints.add(key); apis.append({"endpoint": endpoint, "method": method, "description": description})
    return apis[:20]

def analyze_document(text: str) -> Dict:
    analysis = {"word_count": len(text.split()), "sentence_count": len(safe_sent_tokenize(text)), "key_terms": extract_key_terms(text, top_n=10)}
    return analysis

def extract_key_terms(text: str, top_n: int = 10) -> List[str]:
    stop_words = set(['the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for'])
    words = re.findall(r'\b[a-z]{3,}\b', text.lower())
    words = [w for w in words if w not in stop_words]
    word_freq = Counter(words)
    return [word for word, count in word_freq.most_common(top_n)]

def generate_ai_insight(analysis: Dict, api_count: int, filename: str) -> str:
    """Generates a concise, actionable AI insight based on document analysis."""
    key_terms = analysis.get("key_terms", [])
    
    if api_count > 0 and 'api' in key_terms[:3]:
        api_terms = [t for t in key_terms if t in ['api', 'service', 'integration', 'endpoint']]
        return f"**API Focus:** {api_count} distinct API endpoints were identified. The core integration focuses on **{key_terms[0].upper()}** and **{key_terms[1].upper()}** data exchange."

    elif any(t in ['model', 'mmm', 'reach', 'roi', 'budget'] for t in key_terms[:5]):
        model_terms = [t.upper() for t in key_terms if t in ['model', 'mmm', 'reach', 'roi', 'budget']]
        return f"**Modeling Core:** The document outlines a **{model_terms[0] if model_terms else 'MODEL'}** framework focusing on **{key_terms[0].upper()}** and **{key_terms[1].upper()}** performance analysis."

    return f"**General Topic:** Document centers on **{key_terms[0].upper()}** and **{key_terms[1].upper()}**, providing {analysis.get('word_count', 0)} words of detail."

# -------------------------
# IMPROVED SEMANTIC DIAGRAM GENERATION
# -------------------------

def extract_semantic_relationships(text: str) -> List[Tuple[str, str, str]]:
    """
    Extract semantic relationships between entities using NLP patterns.
    Returns list of (entity1, relationship, entity2) tuples.
    """
    relationships = []
    
    if not nlp:
        return relationships
    
    doc = nlp(text)
    
    # Pattern 1: Subject-Verb-Object relationships
    for sent in doc.sents:
        sent_text = sent.text.lower()
        
        # API-related patterns
        if any(word in sent_text for word in ['api', 'endpoint', 'service']):
            if 'call' in sent_text or 'request' in sent_text:
                relationships.append(('Client', 'calls', 'API'))
            if 'return' in sent_text or 'response' in sent_text:
                relationships.append(('API', 'returns', 'Data'))
            if 'authenticate' in sent_text or 'auth' in sent_text:
                relationships.append(('User', 'authenticates', 'System'))
        
        # Data flow patterns
        if any(word in sent_text for word in ['data', 'information', 'content']):
            if 'flow' in sent_text or 'transfer' in sent_text:
                relationships.append(('Source', 'sends', 'Data'))
            if 'process' in sent_text or 'analyze' in sent_text:
                relationships.append(('System', 'processes', 'Data'))
            if 'store' in sent_text or 'save' in sent_text:
                relationships.append(('Database', 'stores', 'Data'))
        
        # System architecture patterns
        if any(word in sent_text for word in ['system', 'architecture', 'component']):
            if 'connect' in sent_text or 'link' in sent_text:
                relationships.append(('ComponentA', 'connects', 'ComponentB'))
            if 'depend' in sent_text or 'rely' in sent_text:
                relationships.append(('Service', 'depends', 'Database'))
    
    return list(set(relationships))  # Remove duplicates

def detect_document_type(text: str) -> str:
    """Detect the type of document based on content analysis."""
    text_lower = text.lower()
    
    # API Documentation
    api_keywords = ['api', 'endpoint', 'rest', 'graphql', 'get', 'post', 'put', 'delete']
    if any(keyword in text_lower for keyword in api_keywords):
        return "api_documentation"
    
    # System Architecture
    arch_keywords = ['architecture', 'component', 'system', 'microservice', 'deployment']
    if any(keyword in text_lower for keyword in arch_keywords):
        return "system_architecture"
    
    # Data Flow
    data_keywords = ['data flow', 'etl', 'pipeline', 'processing', 'stream']
    if any(keyword in text_lower for keyword in data_keywords):
        return "data_flow"
    
    # Business Process
    business_keywords = ['workflow', 'process', 'business', 'user journey']
    if any(keyword in text_lower for keyword in business_keywords):
        return "business_process"
    
    return "generic"

def generate_api_architecture_diagram(apis: List[Dict], text: str) -> str:
    """Generate a semantic diagram for API documentation."""
    relationships = extract_semantic_relationships(text)
    
    diagram_lines = ["graph TD"]
    
    # Add API endpoints as nodes
    for i, api in enumerate(apis[:8]):  # Limit to 8 APIs for readability
        method = api['method']
        endpoint = api['endpoint'].split('/')[-1] or 'root'
        node_id = f"API{i}"
        label = f"{method}\\n{endpoint}"
        diagram_lines.append(f"    {node_id}[{label}]")
    
    # Add relationships
    for rel in relationships:
        diagram_lines.append(f"    {rel[0]} -->|{rel[1]}| {rel[2]}")
    
    # Add default API flow if no specific relationships found
    if not relationships and apis:
        diagram_lines.append("    Client[Client App] -->|HTTP Request| API0")
        diagram_lines.append("    API0 -->|Process| Database[(Database)]")
        diagram_lines.append("    Database -->|Return Data| API0")
        diagram_lines.append("    API0 -->|JSON Response| Client")
    
    return "\n".join(diagram_lines)

def generate_system_architecture_diagram(text: str) -> str:
    """Generate a semantic diagram for system architecture."""
    relationships = extract_semantic_relationships(text)
    key_terms = extract_key_terms(text, top_n=8)
    
    diagram_lines = ["graph TD"]
    
    # Create nodes from key terms
    nodes = {}
    for i, term in enumerate(key_terms):
        node_id = f"N{i}"
        nodes[term] = node_id
        diagram_lines.append(f"    {node_id}[{term.upper()}]")
    
    # Add relationships
    for rel in relationships:
        # Try to map relationship entities to our key terms
        source = next((term for term in key_terms if term in rel[0].lower()), "Client")
        target = next((term for term in key_terms if term in rel[2].lower()), "System")
        
        if source in nodes and target in nodes:
            diagram_lines.append(f"    {nodes[source]} -->|{rel[1]}| {nodes[target]}")
    
    # Add default system flow
    if not relationships:
        diagram_lines.extend([
            "    Client[Client] -->|Request| Frontend[Frontend]",
            "    Frontend -->|API Call| Backend[Backend]",
            "    Backend -->|Query| Database[(Database)]",
            "    Database -->|Data| Backend",
            "    Backend -->|Response| Frontend",
            "    Frontend -->|Display| Client"
        ])
    
    return "\n".join(diagram_lines)

def generate_data_flow_diagram(text: str) -> str:
    """Generate a semantic diagram for data flow."""
    diagram_lines = ["graph LR"]
    
    # Extract data-related entities
    data_entities = ['source', 'input', 'raw', 'processed', 'output', 'destination']
    found_entities = []
    
    for entity in data_entities:
        if entity in text.lower():
            found_entities.append(entity)
    
    # Build data flow
    if found_entities:
        for i, entity in enumerate(found_entities):
            node_id = entity.upper()
            diagram_lines.append(f"    {node_id}[{entity.title()} Data]")
            
            if i > 0:
                prev_entity = found_entities[i-1].upper()
                diagram_lines.append(f"    {prev_entity} -->|transform| {node_id}")
    else:
        # Default data flow
        diagram_lines.extend([
            "    SOURCE[Data Source] -->|Extract| PROCESS[Processing]",
            "    PROCESS -->|Transform| STORAGE[Data Storage]",
            "    STORAGE -->|Load| DESTINATION[Destination]",
            "    DESTINATION -->|Visualize| DASHBOARD[Dashboard]"
        ])
    
    return "\n".join(diagram_lines)

def generate_business_process_diagram(text: str) -> str:
    """Generate a semantic diagram for business processes."""
    diagram_lines = ["graph TD"]
    
    # Extract process steps from sentences
    sentences = safe_sent_tokenize(text)
    process_steps = []
    
    for sentence in sentences[:6]:  # Limit to 6 steps
        if any(word in sentence.lower() for word in ['first', 'then', 'next', 'after', 'finally']):
            # Extract the main action
            words = sentence.split()[:4]
            step = ' '.join(words).strip('.,;:')
            if step and len(step) > 3:
                process_steps.append(step)
    
    if process_steps:
        for i, step in enumerate(process_steps):
            step_id = f"Step{i}"
            diagram_lines.append(f"    {step_id}[{step}]")
            
            if i > 0:
                prev_step = f"Step{i-1}"
                diagram_lines.append(f"    {prev_step} --> {step_id}")
    else:
        # Default business process
        diagram_lines.extend([
            "    START[Start Process] -->|Initiate| STEP1[Step 1]",
            "    STEP1 -->|Process| STEP2[Step 2]",
            "    STEP2 -->|Validate| STEP3[Step 3]",
            "    STEP3 -->|Complete| END[End Process]"
        ])
    
    return "\n".join(diagram_lines)

def generate_semantic_diagram(text: str) -> str:
    """
    Generate a meaningful semantic diagram based on document content analysis.
    This is the improved version that creates proper semantic diagrams.
    """
    doc_type = detect_document_type(text)
    apis = extract_apis_with_nlp(text)
    key_terms = extract_key_terms(text, top_n=8)
    
    logger.info(f"Detected document type: {doc_type}")
    logger.info(f"Key terms for diagram: {key_terms}")
    
    if doc_type == "api_documentation" and apis:
        logger.info("Generating API Architecture Diagram")
        return generate_api_architecture_diagram(apis, text)
    
    elif doc_type == "system_architecture":
        logger.info("Generating System Architecture Diagram")
        return generate_system_architecture_diagram(text)
    
    elif doc_type == "data_flow":
        logger.info("Generating Data Flow Diagram")
        return generate_data_flow_diagram(text)
    
    elif doc_type == "business_process":
        logger.info("Generating Business Process Diagram")
        return generate_business_process_diagram(text)
    
    else:
        # Generic diagram based on key terms and relationships
        logger.info("Generating Generic Semantic Diagram")
        relationships = extract_semantic_relationships(text)
        
        diagram_lines = ["graph TD"]
        
        # Add key terms as nodes
        for i, term in enumerate(key_terms[:6]):  # Limit to 6 terms
            diagram_lines.append(f"    N{i}[{term.upper()}]")
        
        # Add relationships between key terms
        for rel in relationships:
            # Simple mapping: use first letters of entities
            source_char = rel[0][0].upper() if rel[0] else 'A'
            target_char = rel[2][0].upper() if rel[2] else 'B'
            diagram_lines.append(f"    {source_char} -->|{rel[1]}| {target_char}")
        
        # Add default connections if no relationships found
        if not relationships and len(key_terms) > 1:
            for i in range(len(key_terms) - 1):
                diagram_lines.append(f"    N{i} -->|related to| N{i+1}")
        
        return "\n".join(diagram_lines)

# -------------------------
# ChromaDB Utility Functions
# -------------------------
def get_text_chunks(text: str, chunk_size: int = 500, chunk_overlap: int = 50) -> List[str]:
    """Splits text into chunks using spaCy for sentence-aware splitting."""
    if not nlp: return [text[i:i + chunk_size] for i in range(0, len(text), chunk_size - chunk_overlap)]
    doc = nlp(text); sentences = [sent.text for sent in doc.sents]; chunks = []; current_chunk = ""
    for sentence in sentences:
        if len(current_chunk) + len(sentence) < chunk_size: current_chunk += (" " + sentence).lstrip()
        else:
            if current_chunk: chunks.append(current_chunk)
            overlap_text = current_chunk[-chunk_overlap:] if len(current_chunk) >= chunk_overlap else ""
            current_chunk = (overlap_text + " " + sentence).lstrip()
    if current_chunk: chunks.append(current_chunk)
    return chunks

def embed_and_store_document(filename: str, text: str):
    """Stores the document chunks in ChromaDB Cloud."""
    collection = get_chroma_collection()
    if not collection:
        logger.warning("Chroma DB Collection is not available. Skipping storage.")
        return

    logger.info(f"Chunking and storing {filename}...")
    chunks = get_text_chunks(text)
    ids = [f"{filename}_{i}" for i in range(len(chunks))]
    metadatas = [{"source": filename, "chunk_index": i} for i in range(len(chunks))]
    
    try:
        collection.delete(where={"source": filename})
        collection.add(
            documents=chunks,
            metadatas=metadatas,
            ids=ids
        )
        logger.info(f"Successfully stored {len(chunks)} chunks for {filename} in Chroma DB Cloud.")
    except Exception as e:
        logger.error(f"Failed to add data to Chroma DB Cloud: {e}")

def semantic_search_chromadb(query_text: str, n_results: int = 3) -> List[Dict]:
    """Queries ChromaDB Cloud for relevant chunks based on a query."""
    collection = get_chroma_collection()
    if not collection:
        logger.warning("Chroma DB Collection is not available. Skipping semantic search.")
        return []
    
    try:
        results = collection.query(
            query_texts=[query_text],
            n_results=n_results,
            include=['documents', 'distances', 'metadatas']
        )
        
        rag_results = []
        if results and results.get('documents') and results['documents'][0]:
            for doc, metadata, distance in zip(results['documents'][0], results['metadatas'][0], results['distances'][0]):
                rag_results.append({
                    "source": metadata.get('source', 'N/A'),
                    "chunk": doc,
                    "distance": round(distance, 4) 
                })
        
        logger.info(f"Semantic search found {len(rag_results)} results.")
        return rag_results
        
    except Exception as e:
        logger.error(f"Semantic search failed: {e}")
        return []

# ----------------------------------------------------------------------
# Core Processing Logic (Encapsulated)
# ----------------------------------------------------------------------
def process_document_logic(file_path: str, filename: str, start_time: float) -> Dict:
    file_type = filename.rsplit('.', 1)[1].lower()
    file_hash = calculate_file_hash(file_path)
    
    # 1. Caching Check
    if file_hash in FILE_CACHE:
        elapsed = round(time.time() - start_time, 2)
        logger.info(f"Cache HIT for {filename}. Took {elapsed}s.")
        cached_data = FILE_CACHE[file_hash].copy()
        cached_data['processing_time'] = elapsed
        cached_data['cached'] = True 
        return cached_data

    logger.info(f"Processing: {filename} (Cache MISS)")

    # 2. Text Extraction
    text = extract_text(file_path, file_type)
    if not text: raise ValueError("No text extracted from document.")

    # 3. Chroma DB Store (RAG Indexing)
    embed_and_store_document(filename, text)

    # 4. Semantic Search / RAG Enhancement
    summary_for_rag = summarize_with_nlp(text, num_sentences=2)
    rag_results = semantic_search_chromadb(summary_for_rag, n_results=5)
    
    # 5. Analysis and NLP
    analysis = analyze_document(text)
    summary = summarize_with_nlp(text, num_sentences=5)
    api_details = extract_apis_with_nlp(text)
    
    # USE THE IMPROVED SEMANTIC DIAGRAM GENERATION
    diagram = generate_semantic_diagram(text)
    
    ai_insight = generate_ai_insight(analysis, len(api_details), filename)

    elapsed = round(time.time() - start_time, 2)
    logger.info(f"Completed in {elapsed}s")

    # 6. Construct Result and Cache
    result = {
        "filename": filename,
        "word_count": analysis["word_count"],
        "sentence_count": analysis["sentence_count"],
        "key_terms": analysis["key_terms"],
        "processing_time": elapsed,
        "summary": summary,
        "api_table": api_details,
        "schematic_diagram": diagram,
        "cached": False, 
        "rag_results": rag_results,
        "ai_insight": ai_insight
    }
    FILE_CACHE[file_hash] = result 
    
    return result

# -------------------------
# Flask Routes (remain the same)
# -------------------------
def save_url_content(url: str, upload_folder: str) -> Tuple[str, str]:
    """Downloads content from a URL, saves it locally, and returns the path and filename."""
    if not (url.startswith('http://') or url.startswith('https://')): raise ValueError("Invalid URL provided.")
    try:
        response = requests.get(url, allow_redirects=True, timeout=30)
        response.raise_for_status()
    except requests.exceptions.RequestException as e: raise ValueError(f"Failed to download URL content: {e}")
    content_type = response.headers.get('Content-Type', '').split(';')[0]
    ext_map = {'application/pdf': '.pdf', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx', 'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx', 'text/plain': '.txt'}
    url_path = requests.utils.urlparse(url).path
    file_ext = os.path.splitext(url_path)[1].lower() if os.path.splitext(url_path)[1] else ext_map.get(content_type, '.pdf')
    if file_ext not in ['.pdf', '.docx', '.pptx', '.txt']: raise ValueError(f"Unsupported file type downloaded from URL: {content_type}")
    filename = secure_filename(os.path.basename(url_path) or f"url_upload_{int(time.time())}{file_ext}")
    if not filename.lower().endswith(file_ext): filename = os.path.splitext(filename)[0] + file_ext
    file_path = os.path.join(upload_folder, filename)
    with open(file_path, 'wb') as f: f.write(response.content)
    return file_path, filename


@app.route('/')
def index():
    return render_template("index.html", api_key=API_KEY)

@app.route('/list_files', methods=['GET'])
def list_files():
    try:
        files = [f for f in os.listdir(app.config['UPLOAD_FOLDER']) if allowed_file(f)]
        file_list = []
        for f in files:
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], f)
            size_kb = round(os.path.getsize(file_path) / 1024, 1)
            file_list.append({"filename": f, "size_kb": size_kb})
        return jsonify({"files": file_list})
    except Exception as e: logger.error(f"Error listing files: {e}"); return jsonify({"error": "Failed to list stored files."}), 500


@app.route('/upload', methods=['POST'])
@limiter.limit("5 per minute") 
def upload_file():
    start_time = time.time()
    auth_header = request.headers.get('Authorization')
    if not auth_header or auth_header.split(' ')[-1] != API_KEY:
        return jsonify({"error": "Unauthorized. Missing or invalid API Key."}), 401
    
    file_path = None 
    filename = None
    file_to_delete = False 

    try:
        # Check for URL upload first
        url_link = request.form.get('url_link')
        
        if url_link:
            logger.info(f"Processing URL: {url_link}")
            file_path, filename = save_url_content(url_link, app.config['UPLOAD_FOLDER'])
            file_to_delete = True # URL files are temporary
        
        # Original: Check for file upload
        elif 'file' in request.files:
            file = request.files['file']
            filename = secure_filename(file.filename)
            if filename == '': raise ValueError("No file selected.")
            if not allowed_file(filename): raise ValueError(f"File type not allowed: {filename.rsplit('.', 1)[-1]}")

            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(file_path)
        
        else:
             # FIX 1: Provide clear error if neither file nor URL is present
             return jsonify({"error": "No file or URL link provided. Please submit one."}), 400

        logger.info(f"File saved to: {file_path}")
        
        # Logic to check if the file is over the 10MB soft limit
        if os.path.getsize(file_path) > MAX_FILE_SIZE:
             raise ValueError(f"File size exceeds the 10MB limit ({round(os.path.getsize(file_path)/(1024*1024), 2)}MB).")

        result = process_document_logic(file_path, filename, start_time)
        
        return jsonify(result)

    except Exception as e:
        logger.error(f"Error: {e}")
        return jsonify({"error": f"Processing failed: {str(e)}"}), 500
    
    finally:
        # Clean up URL-downloaded file (if necessary)
        if file_to_delete and file_path and os.path.exists(file_path):
             try: os.remove(file_path)
             except Exception as e: logger.warning(f"Failed to delete temp file: {e}")


@app.route('/analyze_stored', methods=['POST'])
@limiter.limit("5 per minute") 
def analyze_stored_file():
    start_time = time.time()
    auth_header = request.headers.get('Authorization')
    if not auth_header or auth_header.split(' ')[-1] != API_KEY:
        return jsonify({"error": "Unauthorized. Missing or invalid API Key."}), 401

    data = request.json
    filename = data.get('filename')
    
    if not filename: return jsonify({"error": "No filename provided for analysis."}), 400
        
    filename = secure_filename(filename)
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    
    if not os.path.exists(file_path) or not allowed_file(filename):
        return jsonify({"error": f"Stored file '{filename}' not found or is not a processable type."}), 404
        
    try:
        result = process_document_logic(file_path, filename, start_time)
        return jsonify(result)
        
    except Exception as e:
        logger.error(f"Error analyzing stored file: {e}")
        return jsonify({"error": f"Analysis of stored file failed: {str(e)}"}), 500


@app.route('/rag_query', methods=['POST'])
@limiter.limit("10 per minute")
def rag_query_route():
    auth_header = request.headers.get('Authorization')
    if not auth_header or auth_header.split(' ')[-1] != API_KEY:
        return jsonify({"error": "Unauthorized. Missing or invalid API Key."}), 401

    data = request.json
    query = data.get('query')
    
    if not query: return jsonify({"error": "No query text provided."}), 400
    if not get_chroma_collection():
        return jsonify({"error": "Vector database (ChromaDB) is not initialized or failed to load."}), 500
        
    try:
        rag_results = semantic_search_chromadb(query, n_results=5)
        
        llm_answer = "RAG results retrieved. Integrate an LLM to synthesize the final answer from the chunks."
        
        return jsonify({
            "query": query,
            "rag_results": rag_results,
            "llm_answer": llm_answer
        })
        
    except Exception as e:
        logger.error(f"RAG Query failed: {e}")
        return jsonify({"error": f"RAG Query processing failed: {str(e)}"}), 500


@app.route('/download_results', methods=['POST'])
def download_results():
    try:
        data = request.json
        if not data: return jsonify({"error": "No data provided for download."}), 400
        if 'cached' in data: del data['cached']
            
        temp_filename = f"{data.get('filename', 'analysis_results').replace('.', '_')}_{time.time()}.json"
        temp_filepath = os.path.join(app.config['UPLOAD_FOLDER'], temp_filename)
        
        with open(temp_filepath, 'w') as f: json.dump(data, f, indent=4)
        
        return send_file(
            temp_filepath,
            as_attachment=True,
            download_name=f"analysis_results_{data.get('filename', 'document')}.json",
            mimetype='application/json'
        )
    except Exception as e:
        logger.error(f"Download error: {e}")
        return jsonify({"error": "Failed to create download file."}), 500
    finally: pass 

@app.route('/health', methods=['GET'])
def health_check():
    collection = get_chroma_collection()
    return jsonify({
        "status": "healthy",
        "nlp_available": nlp is not None,
        "cache_size": len(FILE_CACHE),
        "stored_files_count": len(os.listdir(app.config['UPLOAD_FOLDER'])),
        "chroma_available": collection is not None,
        "chroma_count": collection.count() if collection else 0
    })

if __name__ == "__main__":
    app.run(debug=True, threaded=True, port=5000)
